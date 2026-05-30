#!/usr/bin/env python3
"""Generate RSM consulting PPT skill HTML and PPTX outputs."""
from __future__ import annotations

import json
import re
import shutil
import sys
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
APP = ROOT / "outputs" / "vqa_template"
CONSULTING_SKILL = Path("/Users/jinkunxiao/.codex/skills/rsm-consulting-ppt-skills/assets/layouts")
SKILL = CONSULTING_SKILL if CONSULTING_SKILL.exists() else ROOT / "PPT skills_副本"
EXPORT = APP / "exports"
ASSETS = EXPORT / "skill_report_assets"

W, H = 20.0, 11.25

NAVY = "061B3A"
DEEP = "0A2F5C"
BLUE = "0099D8"
SKY = "DFF1FC"
RED = "EF4444"
GREEN = "10B981"
AMBER = "F59E0B"
GREY = "2F3A4A"
MID = "667085"
LIGHT = "A6B0BE"
BORDER = "D9E1EA"
LINE = "E5EAF0"
WHITE = "FFFFFF"

FONT = "Noto Sans SC"
FONT_EN = "Inter"


def rgb(hex_color: str) -> RGBColor:
    h = hex_color.strip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def load_data() -> dict:
    text = (APP / "data.js").read_text(encoding="utf-8")
    match = re.search(r"window\.VQA_DATA\s*=\s*(\{.*\});\s*$", text, re.S)
    if not match:
      raise RuntimeError("Cannot parse data.js")
    return json.loads(match.group(1))


DATA = load_data()
RECORDS = DATA["records"]
BANKS = DATA["banks"]
ALIASES = DATA.get("aliases", {})


def resolve_bank(name: str) -> str:
    return ALIASES.get(name, name)


def bank_meta(name: str) -> dict:
    real = resolve_bank(name)
    return next((b for b in BANKS if b["bank"] == real), {}) or next((r for r in RECORDS if r["bank"] == real), {})


def display_bank(name: str) -> str:
    real = resolve_bank(name)
    meta = bank_meta(real)
    if "农村商业银行" in real or real.endswith("银行"):
        return real
    if "农商行" in real:
        return real.replace("农商行", "农村商业银行")
    if "城市商业银行" in (meta.get("type") or ""):
        return f"{real}银行"
    if "农村商业银行" in (meta.get("type") or ""):
        return f"{real}农村商业银行"
    return real


def latest(bank: str, year: int = 2025) -> dict:
    real = resolve_bank(bank)
    return next(r for r in RECORDS if r["bank"] == real and r["year"] == year)


def rows_for(year: int = 2025, bank_type: str | None = None) -> list[dict]:
    return [r for r in RECORDS if r["year"] == year and (bank_type is None or r["type"] == bank_type)]


def avg(rows: list[dict], key: str) -> float | None:
    vals = [r.get(key) for r in rows if isinstance(r.get(key), (int, float))]
    return sum(vals) / len(vals) if vals else None


def rank_desc(rows: list[dict], key: str, bank: str) -> int | None:
    vals = sorted([r for r in rows if isinstance(r.get(key), (int, float))], key=lambda r: r[key], reverse=True)
    for i, r in enumerate(vals, 1):
        if r["bank"] == bank:
            return i
    return None


def rank_asc(rows: list[dict], key: str, bank: str) -> int | None:
    vals = sorted([r for r in rows if isinstance(r.get(key), (int, float))], key=lambda r: r[key])
    for i, r in enumerate(vals, 1):
        if r["bank"] == bank:
            return i
    return None


def fmt(value, digits=2, suffix="%") -> str:
    if value is None:
        return "暂无"
    return f"{value:.{digits}f}{suffix}"


def yuan(value) -> str:
    if value is None:
        return "暂无"
    return f"{value / 10000:.2f} 亿元"


def safe_text(text: str) -> str:
    return str(text).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


@dataclass
class DeckContext:
    target: str = "苏州农商行"
    year: int = 2025

    @property
    def row(self) -> dict:
        return latest(self.target, self.year)

    @property
    def peers(self) -> list[dict]:
        names = ["常熟农商行", "瑞丰农商行", "上海农商行", "无锡农商行", "张家港农商行", "江阴农商行", "紫金农商行", "苏州"]
        return [latest(n, self.year) for n in names if any(r["bank"] == resolve_bank(n) and r["year"] == self.year for r in RECORDS)]

    @property
    def sample(self) -> list[dict]:
        ids = {self.row["bank"], *[r["bank"] for r in self.peers]}
        return [r for r in RECORDS if r["year"] == self.year and r["bank"] in ids]


CTX = DeckContext()


def html_shell(slides: list[str]) -> str:
    chrome = (SKILL / "chrome.css").read_text(encoding="utf-8")
    chrome = chrome.replace(
        "@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800&family=Noto+Sans+SC:wght@300;400;500;600;700;800;900&display=swap');",
        ""
    )
    labels = []
    for idx, slide in enumerate(slides, 1):
        match = re.search('data-screen-label="([^"]+)"', slide)
        labels.append((idx, safe_text(match.group(1) if match else "报告页面")))
    nav = "\n".join(f'<a href="#slide-{idx:02d}">{idx:02d}｜{label}</a>' for idx, label in labels)
    slide_html = "".join(
        slide.replace('class="slide-container"', f'id="slide-{idx:02d}" class="slide-container"', 1)
        for idx, slide in enumerate(slides, 1)
    )
    return f"""<!doctype html>
<html lang="zh-CN">
<head>
<meta charset="utf-8">
<title>{display_bank(CTX.target)}价值创造与经营对标分析汇报</title>
<style>
{chrome}
@page {{ size: 1920px 1080px; margin: 0; }}
body {{ overflow:auto; background:#eef3f8; }}
body, .slide-container {{ font-family:'Noto Sans SC','Source Han Sans CN','PingFang SC','Microsoft YaHei',sans-serif; }}
.num, .en, .metric, .stat {{ font-family:'Inter','Noto Sans SC',sans-serif; }}
.slide-container {{ margin: 24px auto 24px 32px; box-shadow:0 18px 60px rgba(6,27,58,.16); }}
.slide-container.export-render {{ margin: 0; box-shadow: none; }}
.html-report-nav {{
  position:fixed; right:24px; top:24px; bottom:24px; width:260px; overflow:auto;
  background:rgba(255,255,255,.96); border:1px solid #D9E1EA; border-top:6px solid #0099D8;
  box-shadow:0 12px 34px rgba(6,27,58,.14); padding:14px; z-index:1000;
}}
.html-report-nav h3 {{ margin:0 0 10px; color:#061B3A; font-size:16px; }}
.html-report-nav a {{
  display:block; color:#667085; text-decoration:none; font-size:13px; font-weight:800;
  line-height:1.35; padding:8px 9px; border-left:3px solid transparent;
}}
.html-report-nav a:hover {{ color:#061B3A; background:#DFF1FC; border-left-color:#0099D8; }}
@media print {{
  body {{ background:#fff; }}
  .slide-container {{ margin:0; box-shadow:none; page-break-after:always; }}
  .html-report-nav {{ display:none; }}
}}
</style>
</head>
<body>
<nav class="html-report-nav"><h3>页面导航</h3>{nav}</nav>
{slide_html}
</body>
</html>"""


def extract_template_slide(filename: str, replacements: dict[str, str] | None = None) -> str:
    """Use the bundled skill layout directly; only replace content tokens."""
    text = (SKILL / filename).read_text(encoding="utf-8")
    match = re.search(r"(<div class=\"slide-container\"[\s\S]*?</div>)\s*<script", text)
    if not match:
        match = re.search(r"(<div class=\"slide-container\"[\s\S]*?</div>)\s*</body>", text)
    if not match:
        raise RuntimeError(f"Cannot extract slide-container from {filename}")
    slide = match.group(1)
    replacements = replacements or {}
    for old, new in replacements.items():
        slide = slide.replace(old, new)
    return slide


def build_template_html_slides() -> list[str]:
    """Exact HTML slides from PPT skills_副本 layouts."""
    bank = display_bank(CTX.target)
    row = CTX.row
    common = {
        "苏州农商行（苏农银行 · 002966.SZ）": bank,
        "苏农银行": bank,
        "苏农": bank,
        "ROA 0.92%": f"ROA {fmt(row.get('roa'))}",
        "PB 0.62": f"PB {fmt(row.get('pb'), 2, '倍')}",
        "01 / 46": "01 / 10",
        "04 / 46": "03 / 10",
        "11 / 46": "04 / 10",
        "13 / 46": "05 / 10",
    }
    slides = [
        extract_template_slide("layout-cover.html", common),
        toc_slide(2, 10),
        extract_template_slide("layout-kpi-summary.html", common),
        extract_template_slide("layout-comparison-table.html", common),
        extract_template_slide("layout-module-summary.html", common),
        extract_template_slide("layout-comparison-table.html", {**common, "04 / 10": "06 / 10", "同城对照 ·": "盈利质量 ·", "12 项指标对比": "盈利质量指标对比"}),
        extract_template_slide("layout-comparison-table.html", {**common, "04 / 10": "07 / 10", "同城对照 ·": "息差负债 ·", "12 项指标对比": "息差负债指标对比"}),
        extract_template_slide("layout-comparison-table.html", {**common, "04 / 10": "08 / 10", "同城对照 ·": "风险拨备 ·", "12 项指标对比": "风险拨备指标对比"}),
        extract_template_slide("layout-module-summary.html", {**common, "05 / 10": "09 / 10", "行业坐标与同业对标": "资本纪律与市场定价"}),
        extract_template_slide("layout-comparison-table.html", {**common, "04 / 10": "10 / 10", "同城对照 ·": "行动建议 ·", "12 项指标对比": "12 个月行动建议"}),
    ]
    return slides


def write_slide_folder(slides: list[str]) -> Path:
    """Create the skill-style .slides deliverable directory."""
    deck_dir = EXPORT / "rsm_board_skill_report.slides"
    assets_dir = deck_dir / "assets"
    if deck_dir.exists():
        shutil.rmtree(deck_dir)
    assets_dir.mkdir(parents=True, exist_ok=True)
    shutil.copy2(SKILL / "chrome.css", assets_dir / "chrome.css")
    manifest = {
        "title": f"{display_bank(CTX.target)}价值创造与经营对标分析汇报",
        "canvas": {"width": 1920, "height": 1080},
        "slides": []
    }
    for idx, slide in enumerate(slides, 1):
        page = deck_dir / f"{idx:02d}.html"
        html = f"""<!doctype html>
<html lang="zh-CN">
<head>
<meta charset="utf-8">
<title>P{idx:02d}</title>
<link rel="stylesheet" href="assets/chrome.css">
</head>
<body>
{slide}
</body>
</html>"""
        page.write_text(html, encoding="utf-8")
        manifest["slides"].append({"page": idx, "file": page.name})
    (deck_dir / "manifest.json").write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    return deck_dir


def chrome_header(label: str, crumb: str) -> str:
    return f"""<div class="chrome-header"><div class="module-tag"><span class="dot"></span><span class="label">{safe_text(label)}</span></div><div class="breadcrumb">{safe_text(crumb)}</div></div>"""


def chrome_footer(page: int, total: int) -> str:
    return f"""<div class="chrome-footer"><div class="src">数据来源：iFinD · 上市公司年报 · RSM 整理 · 2020—2025 全年口径</div><div class="rsm-brand"><span class="logo-box">RSM</span><span>{display_bank(CTX.target)}价值创造汇报 · 2026.05</span></div><div class="pagenum">{page:02d} / {total:02d}</div></div>"""


def title_block(title: str, subtitle: str) -> str:
    return f"""<div class="page-title-block"><div class="rule"></div><div class="h1">{safe_text(title)}</div><div class="h2">{safe_text(subtitle)}</div></div>"""


def cover_slide(total: int) -> str:
    row = CTX.row
    bank = display_bank(CTX.target)
    peer_text = "、".join(display_bank(r["bank"]) for r in CTX.peers[:5])
    return f"""<div class="slide-container" data-screen-label="01 封面" style="background:linear-gradient(135deg,#FFFFFF 0%,#061B3A 55%,#0099D8 100%);">
<div style="position:absolute;left:1720px;top:60px;width:140px;height:48px;background:white;display:flex;align-items:center;justify-content:center;font-family:Inter,sans-serif;font-weight:800;font-size:24pt;color:#061B3A;letter-spacing:.12em;">RSM</div>
<div style="position:absolute;left:60px;top:80px;width:80px;height:3px;background:#0099D8;"></div>
<div style="position:absolute;left:60px;top:100px;width:600px;height:30px;font-size:13pt;color:#667085;letter-spacing:.3em;font-weight:500;">董　事　会　汇　报</div>
<div style="position:absolute;left:60px;top:200px;width:1500px;font-size:24pt;color:#667085;font-weight:400;letter-spacing:.05em;">{bank}</div>
<div style="position:absolute;left:60px;top:255px;width:1700px;font-size:64pt;color:#061B3A;font-weight:800;line-height:1.1;letter-spacing:-.02em;">价值创造与经营对标分析汇报</div>
<div style="position:absolute;left:60px;top:400px;width:1700px;font-size:24pt;color:#667085;font-weight:300;line-height:1.4;">经营视角与资本市场视角的双重穿透分析</div>
<div style="position:absolute;left:60px;top:500px;width:1800px;height:1px;background:linear-gradient(to right,#0099D8 0%,rgba(255,255,255,.1) 100%);"></div>
<div style="position:absolute;left:60px;top:540px;width:1100px;font-size:16pt;color:#2F3A4A;font-weight:300;line-height:1.7;">
本汇报采用 <strong style="color:#0099D8;font-weight:600;">RSM 银行价值质量评估体系</strong>（VQA）——总资产收益率（ROA）→ 风险调整资本回报率（RAROC）→ 经济增加值（EVA）三层穿透框架，重估 {bank} {CTX.year} 年 <strong style="color:#061B3A;">ROA {fmt(row.get('roa'))}</strong> 的可持续性；并以 PB、市净率和同业指标进行交叉验证。
</div>
<div style="position:absolute;left:1280px;top:540px;width:560px;font-size:13pt;color:#667085;line-height:1.9;border-left:2px solid rgba(255,255,255,.2);padding-left:32px;">
<div style="font-size:11pt;color:#7090A8;letter-spacing:.2em;margin-bottom:14px;">报告元数据</div>
<div><span style="display:inline-block;width:90px;color:#7090A8;">报告期</span>2026 年 5 月</div>
<div><span style="display:inline-block;width:90px;color:#7090A8;">报告口径</span>2020 — 2025</div>
<div><span style="display:inline-block;width:90px;color:#7090A8;">沟通对象</span>董事会 / 董事长</div>
<div><span style="display:inline-block;width:90px;color:#7090A8;">目标银行</span>{bank}</div>
<div><span style="display:inline-block;width:90px;color:#7090A8;">对标样本</span>{peer_text}</div>
<div><span style="display:inline-block;width:90px;color:#7090A8;">数据来源</span>iFinD · 上市公司年报 · RSM 研究</div>
</div>
<div style="position:absolute;left:60px;top:870px;width:1800px;font-size:13pt;color:#7090A8;letter-spacing:.15em;">阅读路径　·　研究口径 → 执行摘要 → VQA 框架 → 行业坐标 → 盈利／中收／息差／风险／资本专题 → 行动方案</div>
<div style="position:absolute;left:60px;top:950px;width:1800px;height:1px;background:rgba(255,255,255,.15);"></div>
<div style="position:absolute;left:60px;top:980px;width:900px;font-size:11pt;color:#7090A8;letter-spacing:.1em;">审计　｜　税务　｜　咨询　　·　　© 2026 RSM 中国 版权所有</div>
<div style="position:absolute;left:1740px;top:980px;width:120px;font-size:13pt;color:#667085;text-align:right;font-family:Inter,sans-serif;">01 / {total:02d}</div>
</div>"""


def toc_slide(page: int, total: int) -> str:
    items = [
        ("01", "摘要与价值质量评估体系（VQA）", "研究口径、执行摘要、VQA框架"),
        ("02", "行业坐标与同业对标", "同区域、同类型与同城样本"),
        ("03", "盈利质量与轻资本收入", "核心营收、手续费、拨备前利润"),
        ("04", "息差防守与负债底盘", "NIM、息差对冲缺口、存款结构"),
        ("05", "风险确认与拨备缓冲", "不良、偏离度、个贷风险、PCR"),
        ("06", "资本纪律与市场定价", "CET1、RWA、PB与行动建议"),
    ]
    rows = "".join(f"""<div style="display:grid;grid-template-columns:86px 1fr;gap:22px;padding:18px 0;border-bottom:1px solid #D9E1EA;"><div style="font-family:Inter,sans-serif;font-size:28pt;font-weight:800;color:#0099D8;">{n}</div><div><div style="font-size:18pt;font-weight:700;color:#061B3A;">{t}</div><div style="font-size:12pt;color:#667085;margin-top:6px;">{d}</div></div></div>""" for n, t, d in items)
    return f"""<div class="slide-container">{chrome_header("P02 · 报告导航", "目录")}
<div style="position:absolute;left:0;top:0;width:560px;height:1080px;background:#F7F8FA;"></div>
<div style="position:absolute;left:60px;top:190px;width:400px;color:#061B3A;"><div style="width:80px;height:4px;background:#0099D8;margin-bottom:26px;"></div><div style="font-size:42pt;font-weight:800;">目录</div><div style="font-size:16pt;color:#667085;line-height:1.7;margin-top:24px;">本报告按 SCQA 节拍展开：先给出经营事实，再解释价值质量差异，最后落到 12 个月管理动作。</div></div>
<div style="position:absolute;left:620px;top:150px;width:1240px;">{rows}</div>
{chrome_footer(page,total)}</div>"""


def kpi_slide(page: int, total: int) -> str:
    row = CTX.row
    sample = CTX.sample
    bank = row["bank"]
    metrics = [
        ("指标 01", "总资产收益率 · ROA", row.get("roa"), "%", rank_desc(sample, "roa", bank), "账面回报水平"),
        ("指标 02", "核心营收同比", row.get("coreRevenueGrowth"), "%", rank_desc(sample, "coreRevenueGrowth", bank), "主业收入修复"),
        ("指标 03", "息差对冲缺口", row.get("nimGapBp"), "bp", rank_asc(sample, "nimGapBp", bank), "资产负债管理"),
        ("指标 04", "个人贷款不良率", row.get("personalLoanNpl") or row.get("retailRiskMax"), "%", rank_asc(sample, "personalLoanNpl", bank), "零售风险暴露"),
    ]
    cards = []
    for i, (tag, label, val, unit, rank, note) in enumerate(metrics):
        color = [BLUE, RED, AMBER, RED][i]
        display = "暂无" if val is None else (f"{val:+.1f}" if unit == "bp" else f"{val:.2f}")
        cards.append(f"""<div style="position:absolute;left:{60+i*460}px;top:316px;width:{445 if i<3 else 420}px;height:340px;background:white;border:1px solid #D9E1EA;border-top:4px solid #{color};padding:24px 26px;">
<div style="font-size:11pt;color:#667085;letter-spacing:.04em;"><span style="display:inline-block;background:#F7F8FA;color:#061B3A;font-size:9pt;padding:2px 8px;border-radius:2px;margin-right:8px;">{tag}</span>{label}</div>
<div style="font-family:Inter,sans-serif;font-size:88pt;font-weight:700;color:#{color};line-height:1;margin-top:12px;">{display}<span style="font-size:36pt;font-weight:500;color:#667085;">{unit}</span></div>
<div style="font-size:11pt;color:#2F3A4A;margin-top:14px;padding-top:12px;border-top:1px dashed #E5EAF0;line-height:1.7;"><span style="background:#E3F2FD;color:#061B3A;font-weight:600;padding:2px 8px;border-radius:2px;">样本排名 {rank or '暂无'} / {len(sample)}</span><br>{note} · 样本均值 {fmt(avg(sample, {'%':'roa','bp':'nimGapBp'}[unit] if label.startswith('息差') else ['roa','coreRevenueGrowth','nimGapBp','personalLoanNpl'][i]), 2 if unit!='bp' else 1, unit)}</div>
</div>""")
    return f"""<div class="slide-container">{chrome_header("① 摘要 + RSM VQA 框架", "P04 · 执行摘要 — 4 项 KPI 全景")}
{title_block(f"{display_bank(CTX.target)} {CTX.year} 年 · 4 项关键绩效指标全景", "阅读路径：行业坐标 → 4 项 KPI → 经营视角与资本市场视角结论 → 战略行动")}
<div style="position:absolute;left:60px;top:236px;width:1800px;height:54px;background:#DFF1FC;border-left:4px solid #0099D8;padding:14px 22px;font-size:14pt;color:#2F3A4A;">行业锚 · 所选样本共 {len(sample)} 家；本页所有排名基于目标银行 + 对标银行的 {CTX.year} 年截面数据。</div>
{''.join(cards)}
<div style="position:absolute;left:60px;top:686px;width:910px;height:282px;background:white;border-left:4px solid #061B3A;padding:24px 28px;"><div style="font-size:10pt;color:#0099D8;font-weight:600;letter-spacing:.12em;margin-bottom:10px;">经营视角结论</div><div style="font-size:14pt;line-height:1.7;color:#2F3A4A;">{display_bank(CTX.target)} {CTX.year} 年营业收入 {yuan(row.get('revenue'))}，归母净利润 {yuan(row.get('netProfit'))}，ROA 为 {fmt(row.get('roa'))}。需要进一步区分账面回报、核心营收修复和风险调整后收益，避免只用利润增速解释经营质量。</div></div>
<div style="position:absolute;left:980px;top:686px;width:880px;height:282px;background:white;border-left:4px solid #0099D8;padding:24px 28px;"><div style="font-size:10pt;color:#0099D8;font-weight:600;letter-spacing:.12em;margin-bottom:10px;">资本市场视角结论</div><div style="font-size:14pt;line-height:1.7;color:#2F3A4A;">PB 为 {fmt(row.get('pb'),2,'倍')}。市净率不应单独解释为价值错配，而应与 ROA、手续费资产比、风险确认和资本余量联读，判断市场定价反映的是价值重估基础还是质量折价。</div></div>
{chrome_footer(page,total)}</div>"""


def module_summary_slide(page: int, total: int, num: str, module: str, conclusion: str, facts: list[str], stats: list[tuple[str, str, str]], next_text: str) -> str:
    cards = "".join(f"""<div style="flex:1;background:white;border:1px solid #D9E1EA;border-top:4px solid #{[BLUE,AMBER,RED][i%3]};padding:22px 24px;"><div style="font-family:Inter,sans-serif;font-size:36pt;font-weight:800;color:#{[BLUE,AMBER,RED][i%3]};line-height:1;margin-bottom:10px;">{i+1:02d}</div><div style="font-size:16pt;font-weight:700;color:#061B3A;margin-bottom:10px;">{safe_text(f.split('：')[0])}</div><div style="font-size:13pt;color:#2F3A4A;line-height:1.65;">{safe_text(f.split('：',1)[1] if '：' in f else f)}</div></div>""" for i, f in enumerate(facts[:3]))
    stat_html = "".join(f"""<div style="flex:1;background:white;border:1px solid #D9E1EA;padding:28px 26px;text-align:center;"><div style="font-size:12pt;color:#667085;margin-bottom:8px;">{safe_text(label)}</div><div style="font-family:Inter,sans-serif;font-size:54pt;font-weight:800;color:#{color};line-height:1;">{safe_text(value)}</div><div style="font-size:12pt;color:#2F3A4A;margin-top:14px;">{safe_text(note)}</div></div>""" for label, value, note, color in stats[:3])
    return f"""<div class="slide-container">{chrome_header(module, f"P{page:02d} · 模块小结")}
<div style="position:absolute;left:0;top:0;width:480px;height:1080px;background:#F7F8FA;"></div>
<div style="position:absolute;left:60px;top:200px;width:400px;z-index:10;"><div style="width:80px;height:4px;background:#0099D8;margin-bottom:24px;"></div><div style="font-size:18pt;color:#667085;letter-spacing:.2em;margin-bottom:12px;">模 块 小 结</div><div style="font-family:Inter,sans-serif;font-size:96pt;font-weight:800;color:#061B3A;line-height:.95;">{num}</div><div style="font-size:28pt;font-weight:700;color:#061B3A;margin-top:18px;line-height:1.3;">{safe_text(module)}</div><div style="font-size:14pt;color:#667085;margin-top:24px;line-height:1.65;">当前样本：目标银行 + 对标银行<br>报告口径：2020—2025<br>输出对象：董事会 / 管理层</div></div>
<div style="position:absolute;left:540px;top:130px;width:1320px;"><div style="font-size:11pt;color:#0099D8;font-weight:600;letter-spacing:.12em;margin-bottom:14px;">模块核心结论</div><div style="font-size:28pt;font-weight:700;color:#061B3A;line-height:1.3;margin-bottom:24px;">{safe_text(conclusion)}</div></div>
<div style="position:absolute;left:540px;top:300px;width:1320px;"><div style="display:flex;gap:20px;">{cards}</div></div>
<div style="position:absolute;left:540px;top:620px;width:1320px;"><div style="font-size:11pt;color:#0099D8;font-weight:600;letter-spacing:.12em;margin-bottom:14px;">本模块核心数据</div><div style="display:flex;gap:20px;">{stat_html}</div></div>
<div style="position:absolute;left:540px;top:900px;width:1320px;height:80px;background:#DFF1FC;border-left:4px solid #0099D8;padding:18px 24px;"><div style="font-size:11pt;color:#0099D8;font-weight:600;letter-spacing:.12em;margin-bottom:6px;">下一模块预告</div><div style="font-size:13.5pt;color:#2F3A4A;line-height:1.5;">{safe_text(next_text)}</div></div>
{chrome_footer(page,total)}</div>"""


def table_slide(page: int, total: int, title: str, subtitle: str, rows: list[tuple[str, str, str, str]]) -> str:
    tr = "".join(
        f"""<tr><td>{safe_text(a)}</td><td>{safe_text(b)}</td><td>{safe_text(c)}</td><td>{safe_text(d)}</td></tr>"""
        for a, b, c, d in rows
    )
    evidence = "；".join(safe_text(f"{a}：{b} / 样本 {c}") for a, b, c, _ in rows[:3])
    implication = safe_text(rows[0][3] if rows else "本页结论需回到底层数据继续复核。")
    if "研究口径" in title:
        question = "本次汇报的结论边界是什么？"
        action = "先冻结样本、指标和年份口径，再进入经营判断。"
    elif "盈利" in title:
        question = "利润表现是否来自可持续主业修复？"
        action = "将核心营收、手续费资产比和 PPOP 纳入季度复盘，避免只用净利润解释经营质量。"
    elif "息差" in title:
        question = "负债端降本能否对冲资产端让价？"
        action = "建立息差对冲缺口月度看板，跟踪资产收益、负债成本和高成本存款压降。"
    elif "风险" in title:
        question = "风险是否已经充分确认，拨备缓冲是否足够？"
        action = "把不良、逾期偏离度、关注率和拨备覆盖率放入同一预警框架。"
    elif "行动" in title:
        question = "董事会应批准哪些可追踪的管理动作？"
        action = "按 0-3、3-6、6-12 个月拆分责任部门、指标阈值和复核频率。"
    else:
        question = "本页结论如何支撑董事会决策？"
        action = "把指标差异转化为责任机制和后续复核事项。"
    return f"""<div class="slide-container" data-screen-label="{page:02d} {safe_text(title)}">{chrome_header("专题证据页", f"P{page:02d} · {title}")}{title_block(title, subtitle)}
<div style="position:absolute;left:60px;top:250px;width:1180px;height:650px;background:white;border:1px solid #D9E1EA;">
  <div style="padding:20px 26px 12px;border-bottom:1px solid #D9E1EA;">
    <div style="font-size:11pt;color:#0099D8;font-weight:700;letter-spacing:.12em;">结论证据矩阵</div>
    <div style="font-size:15pt;color:#061B3A;font-weight:700;margin-top:6px;">每项判断均保留目标值、样本值和董事会阅读含义</div>
  </div>
  <table style="width:100%;border-collapse:collapse;font-size:13.5pt;color:#2F3A4A;">
    <thead><tr style="background:#F7F8FA;color:#061B3A;"><th>指标/动作</th><th>目标银行</th><th>样本/责任</th><th>董事会阅读含义</th></tr></thead>
    <tbody>{tr}</tbody>
  </table>
</div>
<div style="position:absolute;left:1280px;top:250px;width:580px;height:650px;">
  <div style="height:150px;background:#061B3A;color:white;padding:22px 26px;border-top:5px solid #0099D8;">
    <div style="font-size:11pt;color:#DFF1FC;font-weight:700;letter-spacing:.12em;margin-bottom:10px;">本页回答</div>
    <div style="font-size:20pt;font-weight:800;line-height:1.28;">{safe_text(question)}</div>
  </div>
  <div style="height:205px;background:white;border:1px solid #D9E1EA;border-top:0;padding:22px 26px;">
    <div style="font-size:11pt;color:#0099D8;font-weight:700;letter-spacing:.12em;margin-bottom:10px;">关键证据</div>
    <div style="font-size:13.5pt;color:#2F3A4A;line-height:1.65;">{evidence}</div>
  </div>
  <div style="height:205px;background:#F7F8FA;border:1px solid #D9E1EA;border-top:0;padding:22px 26px;">
    <div style="font-size:11pt;color:#0099D8;font-weight:700;letter-spacing:.12em;margin-bottom:10px;">管理含义</div>
    <div style="font-size:13.5pt;color:#2F3A4A;line-height:1.65;">{implication}<br>{safe_text(action)}</div>
  </div>
  <div style="margin-top:14px;background:#DFF1FC;border-left:4px solid #0099D8;padding:14px 18px;font-size:12pt;color:#2F3A4A;line-height:1.5;">
    口径边界：本页所有结论均基于 2020—2025 年年报与 iFinD 数据；强判断需回到事实包复核。
  </div>
</div>
<style>td,th{{border-bottom:1px solid #E5EAF0;padding:15px 18px;text-align:left;vertical-align:top;line-height:1.45;}} th{{font-weight:800;}}</style>
{chrome_footer(page,total)}</div>"""


def build_html_slides() -> list[str]:
    row = CTX.row
    sample = CTX.sample
    total = 10
    slides = [cover_slide(total), toc_slide(2,total), kpi_slide(3,total)]
    slides.append(table_slide(4,total,"研究口径、对标样本与指标定义",f"{display_bank(CTX.target)}｜目标银行 + 对标银行｜2020—2025 年",[
        ("样本边界", display_bank(CTX.target), f"{len(sample)} 家银行", "所有结论先在目标银行与对标组内比较，再与类型均值联读。"),
        ("数据口径", "上市公司年报 / iFinD", "2020—2025", "利润、资产、风险、资本和 PB 指标均来自同一底表。"),
        ("方法论", "VQA 三层穿透", "ROA → RAROC → EVA", "先解释账面回报，再解释风险调整和资本占用。"),
    ]))
    slides.append(module_summary_slide(5,total,"02","行业坐标与同业对标",f"{display_bank(CTX.target)}需要同时解释账面位次和结构位次的差异",[
        f"账面回报：ROA {fmt(row.get('roa'))}，在样本内排名 {rank_desc(sample,'roa',row['bank']) or '暂无'} / {len(sample)}。",
        f"收入结构：手续费资产比 {fmt(row.get('feeAssetRatio'),3)}，用于检验轻资本收入能力。",
        f"风险确认：隐性不良暴露 {fmt(row.get('hiddenNplExposure'))}，需要与拨备覆盖率联读。",
    ],[
        ("ROA 排名", f"{rank_desc(sample,'roa',row['bank']) or '—'} / {len(sample)}", fmt(row.get('roa')), BLUE),
        ("PB 水平", fmt(row.get('pb'),2,"倍"), "资本市场读数", AMBER),
        ("拨备覆盖率", fmt(row.get('provisionCoverage')), "风险缓冲", GREEN),
    ],"下一模块进入盈利质量穿透，重点回答利润增长是否来自主业修复。"))
    slides.append(table_slide(6,total,"盈利质量与轻资本收入穿透", "核心营收、手续费资产比、拨备前利润共同决定利润质量", [
        ("核心营收同比", fmt(row.get("coreRevenueGrowth")), fmt(avg(sample,"coreRevenueGrowth")), "低于样本均值时，应拆分净息收入、手续费和投资收益贡献。"),
        ("手续费资产比", fmt(row.get("feeAssetRatio"),3), fmt(avg(sample,"feeAssetRatio"),3), "该指标用于判断轻资本收入是否足以缓冲息差压力。"),
        ("PPOP 增速", fmt(row.get("ppopGrowth")), fmt(avg(sample,"ppopGrowth")), "拨备前利润能解释利润增长的经营底色。"),
    ]))
    slides.append(table_slide(7,total,"息差防守与负债底盘", "资产端让价与负债端降本需要放在同一张表里阅读", [
        ("净息差", fmt(row.get("nim")), fmt(avg(sample,"nim")), "NIM 反映资产负债定价能力的综合结果。"),
        ("生息资产收益率", fmt(row.get("earningAssetYield")), fmt(avg(sample,"earningAssetYield")), "资产端收益率下行会直接影响收入底盘。"),
        ("计息负债成本率", fmt(row.get("interestLiabilityCost")), fmt(avg(sample,"interestLiabilityCost")), "负债成本下降速度决定息差能否对冲。"),
        ("息差对冲缺口", fmt(row.get("nimGapBp"),1,"bp"), fmt(avg(sample,"nimGapBp"),1,"bp"), "正值表示资产让价快于负债降本。"),
    ]))
    slides.append(table_slide(8,total,"风险确认与拨备缓冲", "不良率、逾期偏离度与拨备覆盖率共同决定利润质量边界", [
        ("不良率", fmt(row.get("npl")), fmt(avg(sample,"npl")), "结果指标，用于判断账面风险水平。"),
        ("逾期/不良偏离度", fmt(row.get("overdueNplDeviation"),2,"倍"), fmt(avg(sample,"overdueNplDeviation"),2,"倍"), "偏离度大于 1 时需要关注风险确认节奏。"),
        ("拨备覆盖率", fmt(row.get("provisionCoverage")), fmt(avg(sample,"provisionCoverage")), "拨备厚度应与隐性风险暴露同时阅读。"),
    ]))
    slides.append(module_summary_slide(9,total,"06","资本纪律与市场定价",f"{display_bank(CTX.target)}的 PB 解释必须回到经营质量、风险确认和资本余量",[
        f"资本余量：核心一级资本缓冲 {fmt(row.get('cet1Buffer'),0,'bp')}，决定扩表和分红空间。",
        f"风险加权资产密度：{fmt(row.get('rwaDensity'))}，用于观察规模扩张的资本消耗。",
        f"市净率：PB {fmt(row.get('pb'),2,'倍')}，需要与 ROA、风险和中收联读。",
    ],[
        ("CET1缓冲", fmt(row.get("cet1Buffer"),0,"bp"), "资本安全边界", GREEN),
        ("RWA密度", fmt(row.get("rwaDensity")), "资本消耗", AMBER),
        ("PB", fmt(row.get("pb"),2,"倍"), "市场定价", BLUE),
    ],"最后进入 12 个月行动地图，把指标改善转化为管理层可执行任务。"))
    slides.append(table_slide(10,total,"12 个月行动建议", "把 VQA 诊断转成董事会可以批准、管理层可以承接的动作", [
        ("0—3 个月", "冻结口径、复核数据、锁定短板指标", "董办 / 财务 / 战略", "先形成可复核底稿，避免汇报口径反复变化。"),
        ("3—6 个月", "拆解息差缺口、存款结构、风险确认节奏", "资产负债 / 风险 / 零售", "把经营压力拆到部门和指标。"),
        ("6—12 个月", "建立核心营收、中收转化、资本消耗和 PB 追踪", "经营管理层", "用季度复盘验证价值质量是否改善。"),
    ]))
    return slides


def add_textbox(slide, x, y, w, h, text, size=12, color=GREY, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    return box


def add_num_textbox(slide, x, y, w, h, text, size=12, color=GREY, bold=False, align=PP_ALIGN.LEFT):
    box = add_textbox(slide, x, y, w, h, text, size, color, bold, align)
    for p in box.text_frame.paragraphs:
        for r in p.runs:
            r.font.name = FONT_EN
    return box


def add_rect(slide, x, y, w, h, color, line=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.color.rgb = rgb(line or color)
    return shape


def ppt_footer(slide, page: int, total: int):
    add_rect(slide, 0, 10.75, 20, 0.5, WHITE, LINE)
    add_textbox(slide, 0.62, 10.92, 8, 0.18, "数据来源：iFinD · 上市公司年报 · RSM 整理", 7, MID)
    add_rect(slide, 9.72, 10.88, 0.65, 0.27, NAVY)
    add_num_textbox(slide, 17.8, 10.9, 1.4, 0.2, f"{page:02d} / {total:02d}", 10, MID, False, PP_ALIGN.RIGHT)


def ppt_header(slide, label: str, crumb: str):
    add_rect(slide, 0.62, 0.25, 0.08, 0.08, BLUE)
    add_textbox(slide, 0.76, 0.21, 7.5, 0.2, label, 10, NAVY, True)
    add_textbox(slide, 12.0, 0.21, 7.3, 0.2, crumb, 10, MID, False, PP_ALIGN.RIGHT)
    add_rect(slide, 0.62, 1.0, 0.62, 0.04, BLUE)


def ppt_title(slide, title: str, subtitle: str):
    add_textbox(slide, 0.62, 1.0, 18.0, 0.92, title, 46, BLUE, True)
    add_textbox(slide, 0.62, 1.98, 18.0, 0.46, subtitle, 22, MID)


def add_panel(slide, x, y, w, h, title, body, accent=BLUE):
    add_rect(slide, x, y, w, h, WHITE, BORDER)
    add_rect(slide, x, y, w, 0.06, accent)
    add_textbox(slide, x + 0.2, y + 0.18, w - 0.4, 0.4, title, 18, NAVY, True)
    add_textbox(slide, x + 0.2, y + 0.74, w - 0.4, h - 0.9, body, 16, GREY)


def build_pptx(path: Path):
    ctx = CTX
    row = ctx.row
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]
    total = 10

    # Cover
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 20, 11.25, WHITE)
    add_rect(slide, 0.62, 0.55, 0.34, 0.06, "8A8F95")
    add_rect(slide, 1.02, 0.55, 0.34, 0.06, GREEN)
    add_rect(slide, 1.42, 0.55, 0.72, 0.06, BLUE)
    add_rect(slide, 17.92, 0.62, 1.46, 0.5, NAVY)
    add_num_textbox(slide, 18.15, 0.68, 1.0, 0.28, "RSM", 24, WHITE, True, PP_ALIGN.CENTER)
    add_textbox(slide, 0.62, 1.08, 3.2, 0.22, "董　事　会　汇　报", 13, BLUE, True)
    add_textbox(slide, 0.62, 2.08, 15.0, 0.42, display_bank(ctx.target), 24, NAVY)
    add_textbox(slide, 0.62, 2.65, 17.0, 1.35, "价值创造与经营对标分析汇报", 64, BLUE, True)
    add_textbox(slide, 0.62, 4.68, 17.0, 0.45, "经营视角与资本市场视角的双重穿透分析", 24, NAVY)
    add_rect(slide, 0.62, 5.55, 11.6, 1.75, "F7F8FA")
    add_rect(slide, 0.62, 5.55, 0.07, 1.75, BLUE)
    add_textbox(slide, 0.84, 5.78, 11.1, 1.2, f"本汇报采用 RSM 银行价值质量评估体系（VQA），重估 {display_bank(ctx.target)} {ctx.year} 年 ROA {fmt(row.get('roa'))} 的可持续性，并以 PB、市净率和同业指标进行交叉验证。", 16, GREY)
    add_rect(slide, 13.38, 5.45, 5.65, 2.75, "F7F8FA")
    add_rect(slide, 13.38, 5.45, 5.65, 0.07, BLUE)
    add_textbox(slide, 13.68, 5.7, 1.2, 0.2, "报告元数据", 11, BLUE, True)
    add_textbox(slide, 13.68, 6.13, 5.05, 1.9, f"报告期：2026 年 5 月\n报告口径：2020—2025\n沟通对象：董事会 / 董事长\n目标银行：{display_bank(ctx.target)}\n对标样本：{len(ctx.sample)} 家银行\n数据来源：iFinD · 上市公司年报", 13, GREY)
    add_textbox(slide, 0.63, 9.1, 17.8, 0.25, "阅读路径 · 研究口径 → 执行摘要 → VQA 框架 → 行业坐标 → 专题分析 → 行动方案", 13, GREY)
    add_num_textbox(slide, 0.62, 10.21, 18.77, 0.25, f"01 / {total:02d}", 11, MID)

    # TOC
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 20, 11.25, WHITE)
    add_rect(slide, 0.62, 0.55, 0.34, 0.06, "8A8F95")
    add_rect(slide, 1.02, 0.55, 0.34, 0.06, GREEN)
    add_rect(slide, 1.42, 0.55, 0.72, 0.06, BLUE)
    add_rect(slide, 0.62, 1.32, 4.86, 7.5, "F7F8FA")
    add_rect(slide, 0.62, 1.32, 0.07, 7.5, BLUE)
    add_textbox(slide, 0.92, 1.65, 3.6, 0.6, "目录", 42, BLUE, True)
    add_textbox(slide, 0.92, 2.45, 4.15, 2.4, "本报告按 SCQA 节拍展开：先给出经营事实，再解释价值质量差异，最后落到 12 个月管理动作。", 16, GREY)
    add_textbox(slide, 6.46, 0.22, 12.93, 0.3, "P02 · 报告导航", 12, NAVY, True)
    for i, title in enumerate(["摘要与价值质量评估体系", "行业坐标与同业对标", "盈利质量与轻资本收入", "息差防守与负债底盘", "风险确认与拨备缓冲", "资本纪律与市场定价"]):
        y = 1.9 + i * 0.84
        add_num_textbox(slide, 6.47, y, 0.7, 0.35, f"{i+1:02d}", 28, BLUE, True)
        add_textbox(slide, 7.4, y + 0.03, 10.8, 0.34, title, 18, NAVY, True)
    ppt_footer(slide, 2, total)

    # Remaining slides: simple faithful panels
    html_slides = build_html_slides()
    titles = [
        ("① 摘要 + RSM VQA 框架", f"{display_bank(ctx.target)} {ctx.year} 年 · 4 项关键绩效指标全景", "阅读路径：行业坐标 → 4 项 KPI → 经营视角与资本市场视角结论 → 战略行动"),
        ("研究口径", "研究口径、对标样本与指标定义", f"{display_bank(ctx.target)}｜目标银行 + 对标银行｜2020—2025 年"),
        ("② 行业坐标与同业对标", f"{display_bank(ctx.target)}需要同时解释账面位次和结构位次的差异", "模块小结页"),
        ("③ 盈利质量穿透", "盈利质量与轻资本收入穿透", "核心营收、手续费资产比、拨备前利润共同决定利润质量"),
        ("④ 息差负债", "息差防守与负债底盘", "资产端让价与负债端降本需要放在同一张表里阅读"),
        ("⑤ 风险拨备", "风险确认与拨备缓冲", "不良率、逾期偏离度与拨备覆盖率共同决定利润质量边界"),
        ("⑥ 资本估值", f"{display_bank(ctx.target)}的 PB 解释必须回到经营质量、风险确认和资本余量", "模块小结页"),
        ("行动建议", "12 个月行动建议", "把 VQA 诊断转成董事会可以批准、管理层可以承接的动作"),
    ]
    metric_rows = [
        [("ROA", fmt(row.get("roa")), "账面回报"), ("核心营收", fmt(row.get("coreRevenueGrowth")), "主业修复"), ("息差缺口", fmt(row.get("nimGapBp"),1,"bp"), "资产负债"), ("PB", fmt(row.get("pb"),2,"倍"), "市场定价")],
        [("样本边界", display_bank(ctx.target), f"{len(ctx.sample)} 家银行"), ("数据口径", "2020—2025", "年报 / iFinD"), ("方法论", "VQA", "ROA → RAROC → EVA")],
        [("ROA排名", f"{rank_desc(ctx.sample,'roa',row['bank']) or '—'} / {len(ctx.sample)}", fmt(row.get("roa"))), ("PB", fmt(row.get("pb"),2,"倍"), "市场读数"), ("拨备覆盖率", fmt(row.get("provisionCoverage")), "风险缓冲")],
        [("核心营收同比", fmt(row.get("coreRevenueGrowth")), fmt(avg(ctx.sample,"coreRevenueGrowth"))), ("手续费资产比", fmt(row.get("feeAssetRatio"),3), fmt(avg(ctx.sample,"feeAssetRatio"),3)), ("PPOP增速", fmt(row.get("ppopGrowth")), fmt(avg(ctx.sample,"ppopGrowth")))],
        [("净息差", fmt(row.get("nim")), fmt(avg(ctx.sample,"nim"))), ("生息资产收益率", fmt(row.get("earningAssetYield")), fmt(avg(ctx.sample,"earningAssetYield"))), ("计息负债成本", fmt(row.get("interestLiabilityCost")), fmt(avg(ctx.sample,"interestLiabilityCost")))],
        [("不良率", fmt(row.get("npl")), fmt(avg(ctx.sample,"npl"))), ("偏离度", fmt(row.get("overdueNplDeviation"),2,"倍"), fmt(avg(ctx.sample,"overdueNplDeviation"),2,"倍")), ("拨备覆盖率", fmt(row.get("provisionCoverage")), fmt(avg(ctx.sample,"provisionCoverage")))],
        [("CET1缓冲", fmt(row.get("cet1Buffer"),0,"bp"), "资本安全"), ("RWA密度", fmt(row.get("rwaDensity")), "资本消耗"), ("PB", fmt(row.get("pb"),2,"倍"), "市场定价")],
        [("0—3个月", "冻结口径", "复核数据"), ("3—6个月", "拆解缺口", "部门承接"), ("6—12个月", "指标追踪", "季度复盘")],
    ]
    for idx, (label, title, subtitle) in enumerate(titles, 3):
        slide = prs.slides.add_slide(blank)
        add_rect(slide, 0, 0, 20, 11.25, "F8F9FB")
        ppt_header(slide, label, f"P{idx:02d}")
        ppt_title(slide, title, subtitle)
        rows = metric_rows[idx - 3]
        for i, (a, b, c) in enumerate(rows[:4]):
            add_panel(slide, 0.62 + i * 4.72, 3.0, 4.35, 2.55, a, f"{b}\n{c}", [BLUE, RED, AMBER, GREEN][i % 4])
        add_panel(slide, 0.62, 6.15, 8.9, 2.65, "经营视角结论", f"{display_bank(ctx.target)}当前判断应先回到指标事实，再解释形成机制。", NAVY)
        add_panel(slide, 10.2, 6.15, 8.9, 2.65, "董事会关注事项", "建议把本页指标纳入后续季度复盘，明确责任部门、目标区间和复核频率。", BLUE)
        ppt_footer(slide, idx, total)

    prs.save(path)


def render_html_slides(html_path: Path, out_dir: Path) -> list[Path]:
    """Render every HTML slide to a 1920×1080 PNG so PPTX matches HTML exactly."""
    try:
        from playwright.sync_api import sync_playwright
    except Exception as exc:  # pragma: no cover - environment dependent
        raise RuntimeError("Python Playwright is required to render skill slides") from exc

    out_dir.mkdir(parents=True, exist_ok=True)
    for old in out_dir.glob("slide_*.png"):
        old.unlink()
    image_paths: list[Path] = []
    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={"width": 1920, "height": 1080}, device_scale_factor=1)
        page.goto(html_path.resolve().as_uri(), wait_until="domcontentloaded")
        page.evaluate(
            """() => {
              document.body.style.background = '#ffffff';
              document.querySelectorAll('.slide-container').forEach((el) => {
                el.classList.add('export-render');
              });
            }"""
        )
        count = page.locator(".slide-container").count()
        for i in range(count):
            path = out_dir / f"slide_{i + 1:02d}.png"
            page.locator(".slide-container").nth(i).screenshot(path=str(path))
            image_paths.append(path)
        browser.close()
    return image_paths


def build_image_pptx(image_paths: list[Path], path: Path) -> None:
    """Create visual-faithful PPTX from rendered slide images."""
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]
    for image_path in image_paths:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(path)


def main() -> None:
    EXPORT.mkdir(parents=True, exist_ok=True)
    ASSETS.mkdir(parents=True, exist_ok=True)
    shutil.copy2(SKILL / "chrome.css", ASSETS / "chrome.css")
    html_path = EXPORT / "rsm_board_skill_report.html"
    pptx_path = EXPORT / "rsm_board_skill_report.pptx"
    editable_pptx_path = EXPORT / "rsm_board_skill_report_editable_draft.pptx"
    html_slides = build_html_slides()
    html_path.write_text(html_shell(html_slides), encoding="utf-8")
    deck_dir = write_slide_folder(html_slides)
    # Keep an editable draft, but use rendered-image PPTX as the main client-facing deck.
    build_pptx(editable_pptx_path)
    try:
        slide_images = render_html_slides(html_path, EXPORT / "skill_report_slides")
        build_image_pptx(slide_images, pptx_path)
    except Exception as exc:
        shutil.copy2(editable_pptx_path, pptx_path)
        print(f"render fallback: {exc}", file=sys.stderr)
    print(html_path)
    print(pptx_path)
    print(editable_pptx_path)
    print(deck_dir)


if __name__ == "__main__":
    main()
